from django.shortcuts import render
from django.http import StreamingHttpResponse, HttpResponseNotFound, FileResponse
from django.core.files.storage import FileSystemStorage
import os
import json
import re
import google.generativeai as genai
from pptx import Presentation
from PIL import Image
import io


# Configure Gemini API key
GEMINI_API_KEY = "AIzaSyC0pMYP1DM4wInATA5596-jdE_v_ZNVYGM"
if not GEMINI_API_KEY:
    raise ValueError("GEMINI_API_KEY environment variable is not set.")

def extract_slide_content(pptx_path):
    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "text": [],
            "images": []
        }
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.has_text_frame and shape.text_frame.text:
                    slide_info["text"].append(shape.text_frame.text.strip())

        # Extract images and save them
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_bytes = shape.image.blob
                image = Image.open(io.BytesIO(image_bytes))
                
                # Save the image temporarily
                img_path = f"temp_slide_{i+1}_image.png"
                image.save(img_path)
                slide_info["images"].append(img_path)
        
        # This line was in the wrong place, it should be outside the inner loops
        slides_data.append(slide_info)
    
    return slides_data



import re

def extract_finding_details(report_text):
    """
    Extracts key details from a single raw report string, ensuring all components are hashable.
    Returns a dictionary of structured data.
    """
    details = {
        'slides': tuple(sorted(re.findall(r'Slide (\d+)', report_text))),
        'type': 'unknown',
        'key_data': []
    }
    
    # Check for numerical inconsistencies
    num_matches = re.findall(r"(\$\d+M|\d+x faster|\d+ mins|\d+ hours|\d+ hours per consultant)", report_text, re.IGNORECASE)
    if num_matches:
        details['type'] = 'numerical'
        details['key_data'] = tuple(sorted(set(num_matches)))
        
    # Check for logical claims
    if 'inconsistent' in report_text.lower() or 'contradiction' in report_text.lower():
        if 'numerical' not in details['type']:
            details['type'] = 'logical'
            
    return details

def create_finding_hash(finding_details):
    """Generates a stable hash for a finding to ensure de-duplication."""
    return hash((details['slides'], details['type'], details['key_data']))

def beautify_output(inconsistencies):
    """
    Consolidates and formats the list of raw inconsistency reports into a single clean text.
    Performs de-duplication based on a structured analysis of each finding.
    """
    unique_findings = {}

    for raw_report in inconsistencies:
        if "no inconsistencies found" in raw_report.lower():
            continue
        
        finding_details = extract_finding_details(raw_report)
        
        # Ensure all data in the hash key is hashable
        if isinstance(finding_details['key_data'], list):
            finding_details['key_data'] = tuple(finding_details['key_data'])
        if isinstance(finding_details['slides'], list):
            finding_details['slides'] = tuple(finding_details['slides'])

        # The hash key itself is a tuple of immutable types
        finding_hash = hash((
            finding_details['slides'],
            finding_details['type'],
            finding_details['key_data']
        ))
        
        if finding_hash not in unique_findings:
            unique_findings[finding_hash] = raw_report

    report_content = ""
    if not unique_findings:
        report_content += "No significant inconsistencies were found in the presentation.\n"
    else:
        cleaned_reports = []
        for report in unique_findings.values():
            cleaned_report = re.sub(r'--- Inconsistency \d+ ---', '', report, flags=re.IGNORECASE | re.DOTALL).strip()
            cleaned_reports.append(cleaned_report)

        report_content += "### Inconsistency Report\n\n"
        for i, report in enumerate(cleaned_reports):
            report_content += f"#### {i+1}. Inconsistency\n\n"
            report_content += report + "\n\n"

    return report_content
def event_stream_generator(request):
    """
    Generator function using an O(n^2) pairwise comparison approach.
    """
    if request.method == 'POST' and request.FILES['pptx_file']:
        try:
            uploaded_file = request.FILES['pptx_file']
            fs = FileSystemStorage()
            file_name = fs.save(uploaded_file.name, uploaded_file)
            file_path = fs.path(file_name)

            genai.configure(api_key=GEMINI_API_KEY)
            model = genai.GenerativeModel('gemini-2.5-flash')

            yield f"data:Analysis started for {uploaded_file.name}...\n\n"
            slides_data = extract_slide_content(file_path)
            yield f"data:Content extracted from {len(slides_data)} slides.\n\n"

            raw_inconsistencies = []

            for i in range(len(slides_data)):
                for j in range(i + 1, len(slides_data)):
                    slide1 = slides_data[i]
                    slide2 = slides_data[j]
                    
                    yield f"data:Comparing Slide {slide1['slide_number']} and Slide {slide2['slide_number']}...\n\n"

                    prompt = f"""
                    Here's a multi-slide presentation. Analyze the following information for factual or logical inconsistencies.
                    Slide {slide1['slide_number']}: Text: {slide1['text']}
                    Slide {slide2['slide_number']}: Text: {slide2['text']}
                    
                    Provide a detailed, structured output referencing the slide numbers and the nature of each issue. Be thorough in your analysis, explaining the conflict, its impact, and the logical gaps. If no inconsistencies are found, state that clearly.
                    """
                    
                    try:
                        response = model.generate_content(prompt)
                        raw_inconsistencies.append(response.text.strip())
                    except Exception as e:
                        raw_inconsistencies.append(f"Error analyzing slides {slide1['slide_number']} and Slide {slide2['slide_number']}: {e}")

            final_beautified_report = beautify_output(raw_inconsistencies)
            with open('final_report.txt', 'w', encoding='utf-8') as f:
                f.write(final_beautified_report)

            yield "data:Consolidating all findings and generating final report...\n\n"
            yield f"data:{json.dumps({'type': 'final_report', 'content': final_beautified_report})}\n\n"

            for slide_info in slides_data:
                for img_path in slide_info['images']:
                    if os.path.exists(img_path):
                        os.remove(img_path)
            if os.path.exists(file_path):
                os.remove(file_path)

            yield "data:END_STREAM\n\n"
        
        except Exception as e:
            yield f"data:An unexpected error occurred: {e}\n\n"

def event_stream(request):
    return StreamingHttpResponse(event_stream_generator(request), content_type='text/event-stream')

def upload_form(request):
    return render(request, 'index.html')

def download_report(request):
    file_path = 'final_report.txt'
    if os.path.exists(file_path):
        response = FileResponse(open(file_path, 'rb'))
        response['Content-Disposition'] = 'attachment; filename="inconsistency_report.txt"'
        return response
    else:
        return HttpResponseNotFound('The report file was not found.')